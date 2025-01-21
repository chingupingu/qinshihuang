from pptx import Presentation
from pptx.util import Inches
import json
from flask import send_file
from io import BytesIO

class PPTXModel:
    def __init__(self, input_path):
        self.input_path = input_path
        self.slides_data = None
        self.headers = None
        self.data = None
        self.json_output = []

    ######################
    ### MAIN FUNCTIONS ###
    ######################

    # these are the functions that are called by the handler

    ### 1. Return PPTX as json ###
    def get_json(self):
        self.read_pptx()
        return self.json_output

    # ### 2. Return json as PPTX ###
    # def get_pptx(self, return_buffer=False):
    #     try:
    #         pptx_buffer = BytesIO()
    #         # Parse the JSON string into a Python list
    #         slides_data = json.loads(self.json_output)
    #         prs = self.write_json_to_pptx(slides_data)
            
    #         # Save the presentation to the buffer
    #         prs.save(pptx_buffer)
    #         pptx_buffer.seek(0)  # Reset buffer position to the start

    #         if return_buffer:
    #             return pptx_buffer

    #         # Return the PPTX as a file
    #         return send_file(
    #             pptx_buffer,
    #             mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    #             as_attachment=False
    #         )

    #     except Exception as e:
    #         raise Exception(f"Error converting JSON to PPTX: {str(e)}")


    ########################
    ### HELPER FUNCTIONS ###
    ########################

    ### PPTX reading ###
    # reads the PPTX and updates self.content
    def read_pptx(self):
        try:
            # Open the PowerPoint file
            prs = Presentation(self.input_path)
            
            # Add debug logging
            total_slides = len(prs.slides)

            ret = []
            
            # Process each slide
            for slide_index, slide in enumerate(prs.slides):
                slide_content = {"title": None, "subtitle": None, "table_data": None}
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # Check for placeholder type
                        if shape.is_placeholder:
                            ph_type = shape.placeholder_format.type
                            if ph_type == 1:  # Title
                                slide_content["title"] = shape.text.strip()
                        else:
                            # Subtitle
                            slide_content["subtitle"] = shape.text.strip()
                    elif shape.has_table:
                        table_data = shape.table
                        table_json = []
                        
                        headers = [cell.text for cell in table_data.rows[0].cells]

                        # Get data from remaining rows
                        for i in range(1, len(table_data.rows)):
                            row_data = [cell.text for cell in table_data.rows[i].cells]
                            table_json.append(row_data)
                        
                        slide_content["table_data"] = {
                            "headers": headers,
                            "data": table_json
                        }
                
                # Process the slide content
                processed_data = self.process_content(slide_content)
                ret.append(processed_data)

            self.json_output = json.dumps(ret)
            return None

        except Exception as e:
            print(f"Exception occurred: {str(e)}")
            return json.dumps({
                "error": f"Failed to read PPTX: {str(e)}"
            })
        
    ### PDF processing ###
    # converts the content to a json object and updates self.processed_content
    def process_content(self, slide_content):
        result = {}

        title = slide_content.get('title')
        result[title] = {}

        # if there is no table data
        if slide_content.get('table_data') == None:
            # Process the subtitle
            subtitle = slide_content.get('subtitle', '')
            if subtitle:
                # Split the subtitle into lines
                lines = subtitle.split('\n')
                sub = lines[0]
                subtitle_dict = { sub: [] }
                for i in range(1, len(lines)):
                    if ':' in lines[i]:
                        key, value = lines[i].split(':', 1)
                        subtitle_dict[sub].append({ key.strip(): value.strip() })
                result[title] = subtitle_dict

        # if there is table data
        else:
            table_data = slide_content.get('table_data')
            if table_data:
                headers = table_data.get('headers', [])
                data = table_data.get('data', [])
                table_list = []
                for row in data:
                    row_dict = {headers[i]: row[i] for i in range(len(headers))}
                    table_list.append(row_dict)
                result[title] = table_list

        return result
            
    # ### Write JSON to PPTX ###
    # def write_json_to_pptx(self, slides_data):
    #     prs = Presentation()

    #     # Iterate over each slide's data
    #     for slide_data in slides_data:
    #         # Add a new slide
    #         slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    #         slide = prs.slides.add_slide(slide_layout)

    #         # Add title if available
    #         for title, content in slide_data.items():
    #             if title:
    #                 title_shape = slide.shapes.title
    #                 title_shape.text = title

    #             # Check if content is a list (table data) or dict (subtitle)
    #             if isinstance(content, list):
    #                 # Add table
    #                 rows, cols = len(content) + 1, len(content[0])
    #                 left = top = Inches(2.0)
    #                 width = Inches(6.0)
    #                 height = Inches(0.8)
    #                 table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    #                 # Set column headings
    #                 for i, header in enumerate(content[0].keys()):
    #                     table.cell(0, i).text = header

    #                 # Add data to table
    #                 for row_idx, row_data in enumerate(content, start=1):
    #                     for col_idx, (key, value) in enumerate(row_data.items()):
    #                         table.cell(row_idx, col_idx).text = value
    #             elif isinstance(content, dict):
    #                 # Add subtitle and bullet points
    #                 left = Inches(1.0)
    #                 top = Inches(2.0)
    #                 width = height = Inches(5.0)
    #                 textbox = slide.shapes.add_textbox(left, top, width, height)
    #                 text_frame = textbox.text_frame

    #                 for subtitle, items in content.items():
    #                     p = text_frame.add_paragraph()
    #                     p.text = subtitle
    #                     for item in items:
    #                         for key, value in item.items():
    #                             bullet = text_frame.add_paragraph()
    #                             bullet.text = f"{key}: {value}"
    #                             bullet.level = 1

    #     return prs